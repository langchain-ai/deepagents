"""Convert a DRBench corpus document to plain text on stdout.

Installed in the task image as `extract-text`. DRBench's corpus is PDF, DOCX, XLSX,
PPTX, and JSONL mailbox exports; the benchmark scores research and synthesis rather
than container-format parsing, so the task provides the same extraction path
upstream's own agent uses instead of leaving the agent to reverse-engineer OOXML.
"""

from __future__ import annotations

import json
import sys
from datetime import datetime, timezone
from pathlib import Path

# Bounds the text handed back for one document so a single pathological file cannot
# flood the agent's context or the trial log.
MAX_OUTPUT_CHARS = 400_000


def _from_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages = []
    for number, page in enumerate(reader.pages, 1):
        pages.append(f"--- page {number} ---\n{page.extract_text() or ''}")
    return "\n\n".join(pages)


def _from_docx(path: Path) -> str:
    import docx

    document = docx.Document(str(path))
    blocks = [paragraph.text for paragraph in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            blocks.append("\t".join(cell.text for cell in row.cells))
    return "\n".join(blocks)


def _from_xlsx(path: Path) -> str:
    import openpyxl

    workbook = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
    sheets = []
    for worksheet in workbook.worksheets:
        rows = [
            "\t".join("" if value is None else str(value) for value in row)
            for row in worksheet.iter_rows(values_only=True)
        ]
        sheets.append(f"--- sheet: {worksheet.title} ---\n" + "\n".join(rows))
    workbook.close()
    return "\n\n".join(sheets)


def _from_pptx(path: Path) -> str:
    from pptx import Presentation

    presentation = Presentation(str(path))
    slides = []
    for number, slide in enumerate(presentation.slides, 1):
        texts = [shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False)]
        slides.append(f"--- slide {number} ---\n" + "\n".join(texts))
    return "\n\n".join(slides)


def _recipients(value: object) -> str:
    """Render a `to`/`cc` field, which is either a list or a bare string."""
    if isinstance(value, list):
        return ", ".join(str(item) for item in value)
    return str(value)


def _render_email(record: dict) -> str:
    """Render a Roundcube mailbox record."""
    sender = record.get("from", "")
    from_name = record.get("from_name")
    if from_name:
        sender = f"{from_name} <{sender}>"
    header = [f"{label}: {value}" for label, value in (
        ("Subject", record.get("subject", "")),
        ("From", sender),
        ("To", _recipients(record.get("to", ""))),
        ("Cc", _recipients(record.get("cc", ""))),
        ("Date", record.get("date", "")),
        ("Folder", record.get("folder", "")),
    ) if value]
    attachments = record.get("attachments")
    if attachments:
        header.append(f"Attachments: {_recipients(attachments)}")
    return "\n".join(header) + f"\n\n{record.get('body', '')}"


def _render_post(post: dict) -> str:
    """Render a Mattermost post."""
    header = [f"{label}: {value}" for label, value in (
        ("Team", post.get("team", "")),
        ("Channel", post.get("channel", "")),
        ("User", post.get("user", "")),
    ) if value]
    created = post.get("create_at")
    if isinstance(created, int):
        # Mattermost stamps posts in milliseconds since the epoch.
        stamp = datetime.fromtimestamp(created / 1000, tz=timezone.utc)
        header.append(f"Date: {stamp.isoformat(sep=' ', timespec='seconds')}")
    return "\n".join(header) + f"\n\n{post.get('message', '')}"


def _render_container(kind: str, container: dict) -> str:
    """Render a Mattermost team or channel definition."""
    name = container.get("display_name") or container.get("name", "")
    header = [f"{kind}: {name}"]
    for label, key in (("Team", "team"), ("Purpose", "purpose"), ("Header", "header")):
        value = container.get(key)
        if value:
            header.append(f"{label}: {value}")
    return "\n".join(header)


def _from_jsonl(path: Path) -> str:
    """Render a JSONL mailbox or chat export as readable messages.

    The corpus uses one tagged-union format for both Roundcube mailboxes and
    Mattermost exports: `email` and `post` records carry the content, while
    `team`, `channel`, and `user` records describe who and what they belong to.
    Directory records are rendered too, since a message only makes sense with the
    channel and people it references.
    """
    blocks = []
    people = []
    for line in path.read_text(encoding="utf-8", errors="replace").splitlines():
        line = line.strip()
        if not line:
            continue
        try:
            record = json.loads(line)
        except json.JSONDecodeError:
            blocks.append(line)
            continue
        if not isinstance(record, dict):
            blocks.append(str(record))
            continue

        kind = record.get("type")
        if kind == "email":
            blocks.append(_render_email(record))
        elif kind == "post" and isinstance(record.get("post"), dict):
            blocks.append(_render_post(record["post"]))
        elif kind in {"team", "channel"} and isinstance(record.get(kind), dict):
            blocks.append(_render_container(kind.capitalize(), record[kind]))
        elif kind == "user":
            # Two user shapes exist: flat, and nested under a `user` key.
            user = record["user"] if isinstance(record.get("user"), dict) else record
            name = " ".join(
                str(user[key]) for key in ("first_name", "last_name") if user.get(key)
            )
            handle = user.get("username", "")
            email = user.get("email", "")
            person = " ".join(part for part in (name, f"({handle})" if handle else "", email) if part)
            if person:
                people.append(person)
        elif kind == "version":
            continue
        else:
            # Unknown record type: keep the raw JSON rather than dropping content.
            blocks.append(json.dumps(record, ensure_ascii=False))

    if people:
        blocks.insert(0, "Directory:\n" + "\n".join(f"- {person}" for person in people))
    return "\n\n---\n\n".join(blocks)


_HANDLERS = {
    ".pdf": _from_pdf,
    ".docx": _from_docx,
    ".xlsx": _from_xlsx,
    ".pptx": _from_pptx,
    ".jsonl": _from_jsonl,
}


def extract(path: Path) -> str:
    """Return `path` as plain text, dispatching on its suffix.

    Args:
        path: Document to convert.

    Returns:
        The document's text, truncated to `MAX_OUTPUT_CHARS`.

    Raises:
        FileNotFoundError: If `path` is not a file.
        ValueError: If `path` has an unsupported suffix.
    """
    if not path.is_file():
        msg = f"not a file: {path}"
        raise FileNotFoundError(msg)
    handler = _HANDLERS.get(path.suffix.lower())
    if handler is None:
        if path.suffix.lower() in {".txt", ".md", ".csv", ".json"}:
            text = path.read_text(encoding="utf-8", errors="replace")
        else:
            msg = f"unsupported file type {path.suffix!r}; supported: {sorted(_HANDLERS)}"
            raise ValueError(msg)
    else:
        text = handler(path)
    if len(text) > MAX_OUTPUT_CHARS:
        text = text[:MAX_OUTPUT_CHARS] + f"\n\n[truncated at {MAX_OUTPUT_CHARS} characters]"
    return text


def main(argv: list[str]) -> int:
    """Print each named document as text. Returns a process exit code."""
    if not argv:
        print("usage: extract-text <file> [<file> ...]", file=sys.stderr)
        return 2
    status = 0
    for name in argv:
        try:
            text = extract(Path(name))
        except (FileNotFoundError, ValueError) as exc:
            print(f"extract-text: {exc}", file=sys.stderr)
            status = 1
            continue
        except Exception as exc:  # noqa: BLE001 - one bad file must not abort the rest
            print(f"extract-text: failed to read {name}: {type(exc).__name__}: {exc}", file=sys.stderr)
            status = 1
            continue
        if len(argv) > 1:
            print(f"===== {name} =====")
        print(text)
    return status


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
